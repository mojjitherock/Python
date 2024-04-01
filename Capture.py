import pyautogui
import time
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

def capture_screen():
    screenshot = pyautogui.screenshot()
    return screenshot

def image_diff(img1, img2, threshold=0.01):
    num_different_pixels = 0
    total_pixels = img1.width * img1.height
    
    for x in range(img1.width):
        for y in range(img1.height):
            if img1.getpixel((x, y)) != img2.getpixel((x, y)):
                num_different_pixels += 1
    
    change_percentage = num_different_pixels / total_pixels
    
    if change_percentage >= threshold:
        diff = Image.new(img1.mode, img1.size)
        for x in range(img1.width):
            for y in range(img1.height):
                if img1.getpixel((x, y)) != img2.getpixel((x, y)):
                    diff.putpixel((x, y), (255, 0, 0))  # Highlight differences in red
        return diff
    else:
        return None

def save_to_ppt(images):
    prs = Presentation()

    for img in images:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        left = top = Inches(1)
        pic = slide.shapes.add_picture(img, left, top)

    prs.save("captured_images.pptx")

def main():
    images = []
    previous_screen = capture_screen()

    while True:
        current_screen = capture_screen()

        diff = image_diff(previous_screen, current_screen, threshold=0.05)  # Adjust the threshold as needed
        if diff is not None:
            timestamp = time.strftime("%Y%m%d-%H%M%S")
            filename = f"screen_{timestamp}.png"
            current_screen.save(filename)
            images.append(filename)
            print(f"Change detected. Saved {filename}")

        previous_screen = current_screen
        time.sleep(90)  # Adjust the interval based on your needs

    save_to_ppt(images)

if __name__ == "__main__":
    main()
